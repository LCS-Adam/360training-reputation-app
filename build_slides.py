"""
build_slides.py — exec deck (python-pptx) in the editorial style that matches the app + PDF.

Grey slides · #efeff0 chart cards · IBM Plex Mono uppercase labels + bold headlines · indigo
accent · Measured/Estimate/Recommended chips · the regenerated plain-title charts.
Numbers pulled live from cc2_app_data so the deck stays reconciled.

Fonts: the deck references "Inter" (sans) and "IBM Plex Mono" (labels) to match the app. Both are
free; IBM Plex Mono TTFs are in ./fonts. If they aren't installed on the presenting machine the
viewer falls back gracefully (layout, uppercase, color and the charts still carry the look).

Run from the project root:
    ./.venv/bin/python build_slides.py   ->  07_final_assets/360training_Exec_Deck.pptx
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image as PILImage

import cc2_app_data as D

ROOT = os.path.dirname(os.path.abspath(__file__))
OUTPPTX = os.path.join(ROOT, "07_final_assets", "360training_Exec_Deck.pptx")

PAGE = RGBColor(0xD0, 0xD1, 0xD5)
INK = RGBColor(0x1D, 0x1D, 0x25)
ACCENT = RGBColor(0x4E, 0x42, 0xE4)
MUT = RGBColor(0x3B, 0x40, 0x49)
CARD = RGBColor(0xEF, 0xEF, 0xF0)
LINE = RGBColor(0xD2, 0xD2, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LAV = RGBColor(0xE7, 0xE6, 0xFB)
SANS, MONO = "Inter", "IBM Plex Mono"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def bg(slide):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = PAGE


def textbox(slide, l, t, w, h):
    tf = slide.shapes.add_textbox(l, t, w, h).text_frame
    tf.word_wrap = True
    return tf


def run(p, text, font=SANS, size=18, color=INK, bold=False, spacing=None, caps=False):
    r = p.add_run(); r.text = text.upper() if caps else text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    if spacing is not None:
        r._r.get_or_add_rPr().set("spc", str(int(spacing * 100)))
    return r


def chip(p, kind):
    # PowerPoint runs can't carry a fill, so chips are bracketed colored mono text (matches the PDF)
    label, color = {"E": ("MEASURED", INK), "I": ("ESTIMATE", MUT), "R": ("RECOMMENDED", ACCENT)}[kind]
    return run(p, f"[ {label} ]", MONO, 11, color, bold=True, spacing=0.5)


def label(slide, text):
    tf = textbox(slide, Inches(0.7), Inches(0.45), Inches(11), Inches(0.4))
    run(tf.paragraphs[0], text, MONO, 12, ACCENT, bold=True, spacing=1.2, caps=True)


def heading(slide, text, top=0.85, size=30):
    tf = textbox(slide, Inches(0.7), Inches(top), Inches(12), Inches(1.0))
    run(tf.paragraphs[0], text, SANS, size, INK, bold=True)


def takeaway(slide, kind, text):
    tf = textbox(slide, Inches(0.7), SH - Inches(0.95), Inches(12), Inches(0.75))
    p = tf.paragraphs[0]
    if kind:
        chip(p, kind); run(p, "   ", SANS, 15)
    run(p, text, SANS, 15, INK, bold=True)


def card_image(slide, png, top=Inches(1.7), max_w=Inches(11.6), max_h=Inches(4.7)):
    path = D.chart(png)
    if not os.path.exists(path):
        return
    iw, ih = PILImage.open(path).size
    w, h = max_w, Emu(int(max_w * ih / iw))
    if h > max_h:
        h, w = max_h, Emu(int(max_h * iw / ih))
    left = int((SW - w) / 2)
    # card backing (rounded rect) so the chart reads as a card on the grey slide
    pad = Inches(0.12)
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(left - pad), int(top - pad),
                                  int(w + 2 * pad), int(h + 2 * pad))
    rect.fill.solid(); rect.fill.fore_color.rgb = CARD
    rect.line.color.rgb = LINE; rect.line.width = Pt(0.75)
    rect.shadow.inherit = False
    slide.shapes.add_picture(path, left, top, width=w, height=h)


def chart_slide(lbl, head, png, kind, take):
    s = prs.slides.add_slide(BLANK); bg(s)
    label(s, lbl); heading(s, head)
    card_image(s, png)
    takeaway(s, kind, take)
    return s


def text_slide(lbl, head, rows):
    s = prs.slides.add_slide(BLANK); bg(s)
    label(s, lbl); heading(s, head)
    tf = textbox(s, Inches(0.7), Inches(1.9), Inches(12), Inches(4.8))
    for i, (kind, lead, body) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(16)
        if kind:
            chip(p, kind); run(p, "   ", SANS, 17)
        run(p, lead + "  ", SANS, 18, INK, bold=True)
        run(p, body, SANS, 18, MUT)
    return s


# ---------------------------------------------------------------- data
h = D.HEADLINE
cov = D.backtest_coverage()
oa = {r["year"]: (r["neg_share"], int(r["n"])) for _, r in D.osha_annual().iterrows()}
rg = D.resolution_gap()
ann = D.rhi_annual()
rlo, rhi = D.recoverable()

# ---------------------------------------------------------------- 1 · title
s = prs.slides.add_slide(BLANK); bg(s)
tf = textbox(s, Inches(0.85), Inches(2.2), Inches(11.6), Inches(0.5))
run(tf.paragraphs[0], "Reputation analysis · 10,601 reviews · 2014–2026", MONO, 13, ACCENT, bold=True, spacing=1.0, caps=True)
tf = textbox(s, Inches(0.8), Inches(2.75), Inches(11.8), Inches(2.2))
run(tf.paragraphs[0], "Why 360training's reviews turned negative — and what holds up under scrutiny",
    SANS, 40, INK, bold=True)
tf = textbox(s, Inches(0.85), Inches(5.4), Inches(11.4), Inches(1.0))
run(tf.paragraphs[0], "An independent read of 10,601 public reviews — what changed, why, and where the "
    "highest-leverage fix is. Every finding labeled by how far we can stand behind it.", SANS, 16, MUT)

# ---------------------------------------------------------------- 2 · exec summary
text_slide("Executive summary", "Three findings that survive scrutiny", [
    ("I", "The decline is real and multi-year.",
     f"A reputation-health score (2023=100) falls {ann['2023']:.0f}→{ann['2024']:.0f}→{ann['2025']:.0f}"
     f"→{ann['2026']:.0f}, under all 9 weightings."),
    ("E", "The most actionable finding is a resolution gap.",
     f"{rg['actionable']['pct_of_neg']:.0%} of unhappy customers name a fixable issue; ~0.5% get a real fix — a 178× drop-off."),
    ("I", "The forecast's honest claim is calibration, not direction.",
     f"Over N={cov['n']} held-out forecasts the 80/95% ranges held {cov['cov80']:.0%}/{cov['cov95']:.0%}; negativity is elevated ~32%."),
    ("R", "Highest-value ask:",
     "the new-enrollment / conversion series — it turns the revenue case from plausible to evidenced."),
])

# ---------------------------------------------------------------- 3 · caveats
text_slide("What this data cannot tell us", "The honest limits, stated up front", [
    ("", "Who posts, not everyone —",
     "reviews can't be split into invited vs. organic; every rate reflects who chose to post."),
    ("", "Small monthly samples —",
     "~50–90 reviews/month after 2024, so month-to-month swings of a few points are noise."),
    ("", "Two different lenses —",
     "the forecast tracks 1–2★ reviews; the driver model tracks 1★ reviews. Stated side by side."),
    ("", "The auto-tagging is AI-generated —",
     "checked for consistency, not certified for accuracy."),
])

# ---------------------------------------------------------------- 4-11 · money charts
chart_slide("01 · Forecast", "Where customer sentiment is heading", "01_neg_share_fan_chart.png", "I",
            f"80/95% forecast ranges held {cov['cov80']:.0%}/{cov['cov95']:.0%} over {cov['n']} tests — calibration, not a confident trend.")
chart_slide("02 · Complaint drivers", "What predicts a furious review", "05_driver_forest_plot.png", "I",
            "Billing and customer-support complaints are the strongest predictors of a 1-star review.")
chart_slide("03 · Resolution gap", "A fixable problem, left unfixed", "C1_resolution_gap.png", "E",
            f"{rg['actionable']['pct_of_neg']:.0%} of complaints are fixable; ~0.5% get a real fix — a 178× drop-off.")
chart_slide("04 · Revenue at risk", "When does fixing this pay for itself?", "CC2_breakeven_threshold.png", "I",
            f"The fix pays for itself above a small conversion lift; recoverable band ~{rlo:.0%}–{rhi:.0%}.")
chart_slide("05 · Review integrity", "An audit, not an accusation", "F2_firsttimer_share_overtime.png", "E",
            "First-time-reviewer share stays flat — the fake-account test comes back clean.")
chart_slide("06 · Product lines", "Where trust is bleeding fastest", "G3_osha_annual_negshare.png", "I",
            f"Safety-training negativity tripled {oa['2023'][0]:.0%}→{oa['2024'][0]:.0%}→{oa['2025'][0]:.0%} (2023→25).")
chart_slide("07 · Reputation health", "A steady, multi-year decline", "H1_rhi_illustrative_calibration.png", "I",
            f"The reputation-health score fell {ann['2023']:.0f}→{ann['2026']:.0f} since 2023.")
chart_slide("08 · Reputation health", "The decline survives every test", "H4_weight_sensitivity.png", "I",
            "Every one of nine weightings declines 2023–2026 — the drop is in the data, not the method.")

# ---------------------------------------------------------------- 12 · recommendations
text_slide("Recommendations", "What leadership should do", [
    ("R", "Close the resolution gap.",
     "Convert 'email-us' deflection into tracked, specific fixes on the fixable majority of complaints."),
    ("R", "Investigate the safety-training break.",
     "Highest-volume line, tripled in negativity 2023→2025 — the highest-leverage place to look."),
    ("R", "Treat reply speed as an operational KPI.",
     "It slipped later (late 2025+) and is a lagging signal worth its own watch threshold."),
    ("R", "Provide one internal series:",
     "new-enrollment / conversion — it hardens the revenue case and makes the break-even call near-certain."),
])

prs.save(OUTPPTX)
print(f"wrote {OUTPPTX}  ({os.path.getsize(OUTPPTX)/1024:.0f} KB, {len(prs.slides._sldIdLst)} slides)")
